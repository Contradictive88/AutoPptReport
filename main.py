import os
import pptx
import pdfplumber
import google.generativeai as genai
from dotenv import load_dotenv

# Load the .env file
load_dotenv()

# Get env variables from the .env file
api_key = os.getenv("GENAI_API_KEY")
surname = os.getenv("SURNAME")

def extract_task_completed(pdf_path):
    task_completed = []

    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    # Check if the row contains the 'Task Completed' column
                    if len(row) > 3 and row[3] != 'Task Completed' and row[3] is not None:
                        task_completed.append(row[3])

    # Filter out empty strings and bullet points
    task_completed = [task.strip('•').strip() for task in task_completed if task.strip()]

    # Join the tasks into a single string
    tasks_string = "\n".join(task_completed)

    return tasks_string

def prompt_to_gemini(tasks):
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-1.0-pro-latest')

    prompt_template = """Listed below are my accomplishments in work. 
                I want you to highlight the important tasks completed for the entire week.
                Please do not duplicate any task that is similar.
                Preserve the bulleted format of it like this and this is also what I want to output to be like this
                The bullets should only be the dotted round not anything else:
                
                • Programmed the webpage for admin side.
                • Added more design towards the front page.

                Here are the tasks for the week:
                """
    
    # Ensure tasks is a single string
    tasks = "\n".join(tasks.splitlines())

    prompt_template += tasks

    response = model.generate_content(prompt_template)
    print(response.text)
    return response.text

def fill_ppt_with_tasks_and_date(pptx_path, tasks, date_range, image_folder):
    prs = pptx.Presentation(pptx_path)
    slide = prs.slides[2]  # Assuming the 3rd slide

    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "[TASK COMPLETED]" in shape.text:
                for paragraph in shape.text_frame.paragraphs:
                    if "[TASK COMPLETED]" in paragraph.text:
                        font = paragraph.runs[0].font
                        color_rgb = font.color.rgb if font.color and hasattr(font.color, 'rgb') else None
                        paragraph.clear()
                        for task in tasks.split('\n'):
                            p = shape.text_frame.add_paragraph()
                            p.text = task
                            p.font.name = font.name
                            p.font.size = font.size
                            p.font.bold = font.bold
                            p.font.italic = font.italic
                            if color_rgb:
                                p.font.color.rgb = color_rgb
                            # Add a blank paragraph to create a space
                            blank_paragraph = shape.text_frame.add_paragraph()
                            blank_paragraph.text = ""
                        break
            if "[DATE-RANGE]" in shape.text:
                for paragraph in shape.text_frame.paragraphs:
                    if "[DATE-RANGE]" in paragraph.text:
                        font = paragraph.runs[0].font
                        color_rgb = font.color.rgb if font.color and hasattr(font.color, 'rgb') else None
                        paragraph.clear()
                        paragraph.text = date_range
                        paragraph.font.name = font.name
                        paragraph.font.size = font.size
                        paragraph.font.bold = font.bold
                        paragraph.font.italic = font.italic
                        if color_rgb:
                            paragraph.font.color.rgb = color_rgb
                        break

    # Replace image placeholders
    image_folder_path = os.path.join(image_folder, date_range)
    if os.path.exists(image_folder_path):
        image_files = [f for f in os.listdir(image_folder_path) if os.path.isfile(os.path.join(image_folder_path, f))]
        image_files.sort()  # Ensure consistent order

        image_placeholders = [shape for shape in slide.shapes if shape.shape_type == 13]  # Placeholder shape type
        for i, image_placeholder in enumerate(image_placeholders):
            if i < len(image_files):
                image_path = os.path.join(image_folder_path, image_files[i])
                left = image_placeholder.left
                top = image_placeholder.top
                width = image_placeholder.width
                height = image_placeholder.height
                slide.shapes.add_picture(image_path, left, top, width, height)
                slide.shapes._spTree.remove(image_placeholder._element)  # Remove the old placeholder
    else:
        print(f"No images found in folder for date range: {date_range}")

    # Define the output folder and file name
    output_folder = r"E:\Python Programming\AutoPptReport\output_folder"
    output_path = os.path.join(output_folder, f"{surname} - {date_range} - Weekly Report.pptx")

    # Ensure the output folder exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    prs.save(output_path)

input_folder = r'E:\Python Programming\AutoPptReport\input_folder'
pptx_template_path = r'E:\Python Programming\AutoPptReport\weekly_report_ppt_template\Weekly Report Template.pptx'

# Get the list of PDF files in the input folder
pdf_files = [f for f in os.listdir(input_folder) if f.endswith('.pdf')]

# Determine the number of files to process
num_files_to_process = input("Enter the number of files you want to process (or 'all' for all files): ")

if num_files_to_process.lower() == 'all':
    num_files_to_process = len(pdf_files)
else:
    num_files_to_process = int(num_files_to_process)

# Process the specified number of files
for pdf_file in pdf_files[:num_files_to_process]:
    # Extract the date range from the file name
    date_range = pdf_file.split(' - ')[1].replace(' - Weekly Accomplishment Report.pdf', '')
    
    # Generate the full path to the PDF file
    pdf_path = os.path.join(input_folder, pdf_file)
    
    # Extract tasks and generate formatted tasks
    tasks = extract_task_completed(pdf_path)
    formatted_tasks = prompt_to_gemini(tasks)
    
    # Fill the PowerPoint with tasks, date range, and images
    fill_ppt_with_tasks_and_date(pptx_template_path, formatted_tasks, date_range, input_folder)